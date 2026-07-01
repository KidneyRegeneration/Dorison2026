#!/usr/bin/env python3
"""
Generate a PowerPoint presentation from a folder of QC GIFs.
Each slide contains GLOM, Marker A, and Marker B QC GIFs for a single sample.

Assumes filenames contain a timestamp in format YYYY-MM-DD_hh.mm.ss
and follow pattern:
    <MARKER>_QC_<...>_YYYY-MM-DD_hh.mm.ss.gif

Usage:
    python generate_qc_powerpoint.py \
        -i /path/to/qc_folder \
        -o output.pptx \
        --marker-a DAPI \
        --marker-b CD45
"""

import argparse
import os
import re
from pathlib import Path
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import tempfile


def parse_arguments():
    parser = argparse.ArgumentParser(description="Generate QC PowerPoint report")
    parser.add_argument("-i", "--input", required=True, help="Folder with QC GIFs")
    parser.add_argument("-o", "--output", required=True, help="Output PowerPoint file")
    parser.add_argument("--marker-a", required=True, help="Marker A name")
    parser.add_argument("--marker-b", required=True, help="Marker B name")
    return parser.parse_args()



from PIL import Image
import tempfile

def extract_and_compress_frames(gif_path, max_size=(500, 500), max_frames=3):
    """
    Extract up to `max_frames` frames from a GIF,
    resize/compress them, and return list of temp file paths.
    """
    frames = []
    try:
        img = Image.open(gif_path)

        for i in range(max_frames):
            try:
                img.seek(i)
            except EOFError:
                break

            frame = img.copy()

            # Resize
            frame.thumbnail(max_size, Image.LANCZOS)

            # Convert to RGB (safer for PowerPoint)
            frame = frame.convert("RGB")

            # Save compressed JPEG
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
            frame.save(tmp.name, "JPEG", quality=70, optimize=True)

            frames.append(tmp.name)

    except Exception as e:
        print(f"Failed to process {gif_path}: {e}")

    return frames

def extract_timestamp(filename):
    """Extract timestamp YYYY-MM-DD_hh.mm.ss from filename."""
    match = re.search(r"(\d{4}-\d{2}-\d{2}_\d{2}\.\d{2}\.\d{2})", filename)
    return match.group(1) if match else None


def organize_by_sample(folder, marker_a, marker_b):
    """Group files by timestamp sample."""
    samples = defaultdict(dict)

    for filepath in os.listdir(folder):
        if not filepath.endswith(".gif"):
            continue

        timestamp = extract_timestamp(filepath)
        if not timestamp:
            continue

        fpath = os.path.join(folder, filepath)
        fname_upper = fpath.upper()
        filename_root = filepath[8:-4]
        if "GLOM" in fname_upper:
            samples[timestamp]["glom"] = fpath
            samples[timestamp]["sample_name"] = filename_root
        elif marker_a.upper() in fname_upper:
            samples[timestamp]["markerA"] = fpath
        elif marker_b.upper() in fname_upper:
            samples[timestamp]["markerB"] = fpath

    return samples


def create_presentation(samples, output_path, marker_a, marker_b):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "QC Report"
    slide.placeholders[1].text = f"GLOM, {marker_a}, {marker_b}"

    for sample_id in sorted(samples.keys()):
        files = samples[sample_id]

        if not all(k in files for k in ["glom", "markerA", "markerB"]):
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = title_box.text_frame
        tf.text = f"Sample: {files['sample_name']}"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        
        # Extract frames for each channel
        glom_frames = extract_and_compress_frames(files["glom"])
        a_frames = extract_and_compress_frames(files["markerA"])
        b_frames = extract_and_compress_frames(files["markerB"])

        # Ensure equal number of frames
        num_frames = min(len(glom_frames), len(a_frames), len(b_frames), 3)

        # Layout settings
        img_w = Inches(2.5)
        img_h = Inches(1.8)

        x_positions = [Inches(0.5), Inches(3.25), Inches(6.0)]
        y_start = Inches(1.0)
        y_gap = Inches(2.0)

        columns = [glom_frames, a_frames, b_frames]

        # Place images (3 rows × 3 columns)
        for row in range(num_frames):
            y = y_start + row * y_gap

            for col, x in enumerate(x_positions):
                try:
                    img_path = columns[col][row]
                    slide.shapes.add_picture(img_path, x, y, width=img_w, height=img_h)
                except Exception as e:
                    print(f"Error placing image: {e}")
       
        # Labels
        label_y = Inches(0.68)

        for x, label in zip(x_positions, ["GLOM", marker_a, marker_b]):
            box = slide.shapes.add_textbox(x, label_y, Inches(2.5), Inches(0.3))
            tf = box.text_frame
            tf.text = label
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    print(f"Saved: {output_path}")


def main():
    args = parse_arguments()

    if not os.path.isdir(args.input):
        print(f"Error: Input folder not found: {args.input}")
        return 1

    samples = organize_by_sample(args.input, args.marker_a, args.marker_b)

    if not samples:
        print("No valid samples found.")
        return 1

    create_presentation(samples, args.output, args.marker_a, args.marker_b)
    return 0


if __name__ == "__main__":
    exit(main())
